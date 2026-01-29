import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

def hex_to_rgb(hex_str):
    """Convert hex string (e.g. 'FF0000') to RGBColor object."""
    hex_str = hex_str.lstrip('#')
    return RGBColor.from_string(hex_str)

def process_text_frame(text_frame, color_obj):
    """Apply color to all runs in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color_obj

def process_shape(shape, color_obj):
    """Recursively process shapes to find text and apply color."""
    # 1. Handle Groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                process_shape(child, color_obj)
        return

    # 2. Handle Tables
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame"):
                    process_text_frame(cell.text_frame, color_obj)
        return

    # 3. Handle Normal Text Shapes
    if hasattr(shape, "text_frame"):
        process_text_frame(shape.text_frame, color_obj)

def main():
    parser = argparse.ArgumentParser(description="Change font color of all text in a PowerPoint file.")
    parser.add_argument("input_path", help="Path to the input .pptx file")
    parser.add_argument("color", help="Hex color code (e.g. 000000 for black, FF0000 for red)")
    parser.add_argument("--output", help="Path to output file (default: input_colored.pptx)")
    
    args = parser.parse_args()
    
    input_path = Path(args.input_path)
    if not input_path.exists():
        print(f"Error: File '{input_path}' not found.")
        sys.exit(1)
        
    output_path = args.output
    if not output_path:
        output_path = input_path.with_name(f"{input_path.stem}_colored{input_path.suffix}")
    
    print(f"Processing '{input_path}'...")
    print(f"Target Color: #{args.color}")
    
    try:
        prs = Presentation(input_path)
        color_obj = hex_to_rgb(args.color)
        
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                process_shape(shape, color_obj)
            count += 1
            print(f"Processed Slide {count}/{len(prs.slides)}", end='\r')
            
        print(f"\nSaving to '{output_path}'...")
        prs.save(output_path)
        print("Done!")
        
    except Exception as e:
        print(f"An error occurred: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
