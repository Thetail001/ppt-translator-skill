"""PowerPoint translation pipeline utilities."""
from __future__ import annotations

import json
import re
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Optional, Any, Dict, List
from xml.dom import minidom

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt, Emu

from .translation import TranslationService


def get_alignment_value(alignment_str: str | None):
    """Convert alignment string to PP_ALIGN enum value."""
    alignment_map = {
        "PP_ALIGN.CENTER": PP_ALIGN.CENTER,
        "PP_ALIGN.LEFT": PP_ALIGN.LEFT,
        "PP_ALIGN.RIGHT": PP_ALIGN.RIGHT,
        "PP_ALIGN.JUSTIFY": PP_ALIGN.JUSTIFY,
        "None": None,
        None: None,
    }
    return alignment_map.get(alignment_str)


def get_vertical_anchor_value(anchor_str: str | None):
    """Convert vertical anchor string to MSO_ANCHOR enum value."""
    anchor_map = {
        "TOP (1)": MSO_ANCHOR.TOP,
        "MIDDLE (3)": MSO_ANCHOR.MIDDLE,
        "BOTTOM (2)": MSO_ANCHOR.BOTTOM,
        "MSO_ANCHOR.TOP": MSO_ANCHOR.TOP,
        "MSO_ANCHOR.MIDDLE": MSO_ANCHOR.MIDDLE,
        "MSO_ANCHOR.BOTTOM": MSO_ANCHOR.BOTTOM,
        "None": None,
        None: None,
    }
    return anchor_map.get(anchor_str)


def get_text_frame_properties(text_frame) -> Dict[str, Any]:
    """Extract properties and content from a text frame (shape or cell)."""
    data: Dict[str, Any] = {
        "paragraphs": []
    }
    
    # Extract global text frame properties if needed (e.g. margins), 
    # but most styling is on paragraphs/runs.

    for paragraph in text_frame.paragraphs:
        p_data = {
            "runs": [], # Store individual runs to preserve mixed formatting
            "alignment": None,
            "line_spacing": None,
            "space_before": None,
            "space_after": None,
            "level": paragraph.level
        }
        
        # Paragraph formatting
        if getattr(paragraph, "line_spacing", None) is not None:
            p_data["line_spacing"] = paragraph.line_spacing
        if getattr(paragraph, "space_before", None) is not None:
            p_data["space_before"] = paragraph.space_before
        if getattr(paragraph, "space_after", None) is not None:
            p_data["space_after"] = paragraph.space_after
        if getattr(paragraph, "alignment", None) is not None:
            p_data["alignment"] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None

        # Extract info for each run
        for run in paragraph.runs:
            run_data = {
                "text": run.text,
                "font_size": None,
                "font_name": None,
                "bold": None,
                "italic": None,
                "font_color": None,
            }
            
            if getattr(run.font, "size", None) is not None:
                run_data["font_size"] = run.font.size.pt
            if getattr(run.font, "name", None):
                run_data["font_name"] = run.font.name
            if hasattr(run.font, "bold"):
                run_data["bold"] = run.font.bold
            if hasattr(run.font, "italic"):
                run_data["italic"] = run.font.italic
            
            # Safe color extraction
            if getattr(run.font, "color", None):
                try:
                    # check type if possible, or just try-except rgb access
                    if hasattr(run.font.color, "rgb") and run.font.color.rgb is not None:
                        run_data["font_color"] = str(run.font.color.rgb)
                    # Note: We currently skip Theme Colors to avoid complexity, 
                    # meaning theme-colored text will revert to default theme color, 
                    # which is usually correct for the new PPT.
                except Exception:
                    pass

            p_data["runs"].append(run_data)
            
        data["paragraphs"].append(p_data)
        
    return data


def remove_control_characters(s):
    """Remove non-printable control characters that can corrupt PPTX."""
    if not s:
        return s
    # Keep newlines (\n, \r) and tabs (\t), remove others < 32
    return "".join(c for c in s if c >= ' ' or c in '\n\r\t')


def apply_text_frame_properties(text_frame, data: Dict[str, Any], font_scale: float = 0.7):
    """Apply properties and content to a text frame."""
    if not data.get("paragraphs"):
        return

    # Clear existing content. 
    # text_frame.clear() removes all paragraphs but leaves one empty one.
    text_frame.clear() 

    for idx, p_data in enumerate(data["paragraphs"]):
        # Reuse the first empty paragraph for the first item, add new ones for others
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
            
        if "level" in p_data:
            paragraph.level = p_data["level"]

        # Apply paragraph formatting
        # DEBUG: Commenting out to prevent corruption
        # if p_data.get("alignment"):
        #     paragraph.alignment = get_alignment_value(p_data["alignment"])
        # if p_data.get("line_spacing"):
        #     paragraph.line_spacing = p_data["line_spacing"]
        # if p_data.get("space_before"):
        #     paragraph.space_before = p_data["space_before"]
        # if p_data.get("space_after"):
        #     paragraph.space_after = p_data["space_after"]

        # Reconstruct runs
        runs_data = p_data.get("runs", [])
        
        # Backward compatibility
        if not runs_data and "text" in p_data:
             runs_data = [{
                 "text": p_data["text"],
                 "font_size": p_data.get("font_size"),
                 "font_name": p_data.get("font_name"),
                 "bold": p_data.get("bold"),
                 "italic": p_data.get("italic"),
                 "font_color": p_data.get("font_color")
             }]

        for run_info in runs_data:
            text_content = run_info.get("text", "")
            if not text_content:
                continue
                
            run = paragraph.add_run()
            # SANITIZE: This is critical for preventing file corruption
            run.text = remove_control_characters(text_content)

            # Apply font formatting
            # DEBUG: Commenting out formatting to isolate corruption source
            # if run_info.get("font_size"):
            #     adjusted_size = run_info["font_size"] * font_scale
            #     run.font.size = Pt(adjusted_size)
            
            # if run_info.get("font_name"):
            #     run.font.name = run_info["font_name"]
            
            # if run_info.get("font_color"):
            #     try:
            #         run.font.color.rgb = RGBColor.from_string(run_info["font_color"])
            #     except Exception:
            #         pass 

            # # Only assign if not None to avoid clearing inheritance
            # if run_info.get("bold") is not None:
            #     run.font.bold = run_info["bold"]
            # if run_info.get("italic") is not None:
            #     run.font.italic = run_info["italic"]


def get_shape_properties(shape):
    """Extract text shape properties."""
    shape_data = {
        "width": shape.width,
        "height": shape.height,
        "left": shape.left,
        "top": shape.top,
        "text_content": {} # container for text frame data
    }
    
    if hasattr(shape, "text_frame"):
        shape_data["text_content"] = get_text_frame_properties(shape.text_frame)
        
    return shape_data


def apply_shape_properties(shape, shape_data):
    """Apply saved properties to a shape."""
    try:
        # SKIP GEOMETRY RESTORATION to prevent file corruption.
        # Translating text changes its length; forcing original dimensions 
        # often breaks the layout engine or corrupts the file.
        # We only want to restore the text content and its internal formatting.
        
        # if "width" in shape_data: shape.width = Emu(shape_data["width"])
        # if "height" in shape_data: shape.height = Emu(shape_data["height"])
        # if "left" in shape_data: shape.left = Emu(shape_data["left"])
        # if "top" in shape_data: shape.top = Emu(shape_data["top"])
        
        if hasattr(shape, "text_frame") and "text_content" in shape_data:
            apply_text_frame_properties(shape.text_frame, shape_data["text_content"], font_scale=0.7)
            
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error applying shape properties: {exc}")


def get_table_properties(table):
    """Extract table properties."""
    table_data = {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "cells": [],
    }
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                # Skip geometry margins extraction if we don't plan to enforce them strictly?
                # But margins are part of text frame props usually.
                # Actually cell.margin_left is a property of the cell.
                "margin_left": cell.margin_left,
                "margin_right": cell.margin_right,
                "margin_top": cell.margin_top,
                "margin_bottom": cell.margin_bottom,
                "vertical_anchor": str(cell.vertical_anchor) if cell.vertical_anchor else None,
                "text_content": get_text_frame_properties(cell.text_frame)
            }
            row_data.append(cell_data)
        table_data["cells"].append(row_data)
    return table_data


def apply_table_properties(table, table_data):
    """Apply saved table properties."""
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                cell_data = table_data["cells"][row_idx][col_idx]
                
                # We can safely restore margins as they are internal padding
                try:
                    cell.margin_left = Emu(cell_data["margin_left"])
                    cell.margin_right = Emu(cell_data["margin_right"])
                    cell.margin_top = Emu(cell_data["margin_top"])
                    cell.margin_bottom = Emu(cell_data["margin_bottom"])
                except (ValueError, TypeError):
                    pass
                
                if cell_data.get("vertical_anchor"):
                    cell.vertical_anchor = get_vertical_anchor_value(cell_data["vertical_anchor"])
                
                if "text_content" in cell_data:
                     apply_text_frame_properties(cell.text_frame, cell_data["text_content"], font_scale=0.8)
                     
            except Exception as exc:  # pragma: no cover - best effort logging
                print(f"Error setting cell properties: {exc}")


def serialize_runs_to_tagged_text(runs_data: List[Dict[str, Any]]) -> str:
    """Convert runs to tagged text for translation (e.g., <r0>Hello</r0><r1>World</r1>)."""
    tagged_parts = []
    for idx, run in enumerate(runs_data):
        text = run.get("text", "")
        if text:
            # Simple XML escaping might be needed if text contains < or >
            # But deepseek usually handles mixed content well.
            # Let's do basic escaping to be safe.
            safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            tagged_parts.append(f"<r{idx}>{safe_text}</r{idx}>")
    return "".join(tagged_parts)


def parse_tagged_text_to_runs(tagged_text: str, original_runs: List[Dict[str, Any]]) -> None:
    """Parse translated tagged text and update original runs."""
    # Pattern to match <rN>content</rN>
    # Uses non-greedy match for content
    pattern = re.compile(r"<r(\d+)>(.*?)</r\1>", re.DOTALL)
    
    matches = pattern.findall(tagged_text)
    
    # If no tags found but we expected tags (and input wasn't empty), it's a failure.
    # However, if original had no text, matches will be empty, which is fine.
    
    found_ids = set()
    
    if matches:
        # Clear all runs first to avoid ghost text from missing tags
        for run in original_runs:
            run["text"] = ""
            
        for run_id_str, content in matches:
            try:
                run_id = int(run_id_str)
                if 0 <= run_id < len(original_runs):
                    # Unescape basic entities
                    text = content.replace("&lt;", "<").replace("&gt;", ">").replace("&amp;", "&")
                    
                    original_runs[run_id]["text"] = text
                    found_ids.add(run_id)
            except ValueError:
                continue
    else:
        # Fallback: AI likely stripped tags.
        print(f"[DEBUG] WARNING: No tags found in translated text: {tagged_text[:50]}...")
        # Put everything in the first run if it exists.
        if original_runs:
            # Strip potential leftover broken tags if any, or just take raw
            clean_text = re.sub(r"</?r\d+>", "", tagged_text).strip()
            if clean_text:
                # Clear all first
                for run in original_runs:
                    run["text"] = ""
                # Assign to first run
                original_runs[0]["text"] = clean_text


def process_shape(
    shape,
    shape_index_str: str,
    parent_element: ET.Element,
    translation_tasks: List[tuple],
    deferred_writes: List[tuple],
):
    """Recursively process a shape or group, collecting translation tasks and deferred XML writes."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if hasattr(shape, "shapes"):
            for child_idx, child_shape in enumerate(shape.shapes):
                child_index_str = f"{shape_index_str}:{child_idx}"
                process_shape(child_shape, child_index_str, parent_element, translation_tasks, deferred_writes)
        return

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table_element = ET.SubElement(parent_element, "table_element")
        table_element.set("shape_index", shape_index_str)
        table_data = get_table_properties(shape.table)
        
        # Iterate through all cells and their paragraphs
        for row in table_data["cells"]:
            for paragraph in row["text_content"]["paragraphs"]:
                if "runs" in paragraph and paragraph["runs"]:
                    tagged_text = serialize_runs_to_tagged_text(paragraph["runs"])
                    if tagged_text.strip():
                        # Defer translation
                        translation_tasks.append((paragraph, tagged_text))
                elif "text" in paragraph and paragraph["text"].strip():
                    translation_tasks.append((paragraph, paragraph["text"]))

        props_element = ET.SubElement(table_element, "properties")
        # Defer writing JSON to XML until after translation updates table_data
        deferred_writes.append((props_element, table_data))
        
    elif hasattr(shape, "text_frame"): 
        text_element = ET.SubElement(parent_element, "text_element")
        text_element.set("shape_index", shape_index_str)
        shape_data = get_shape_properties(shape)
        
        if "text_content" in shape_data:
            for paragraph in shape_data["text_content"]["paragraphs"]:
                 if "runs" in paragraph and paragraph["runs"]:
                    tagged_text = serialize_runs_to_tagged_text(paragraph["runs"])
                    if tagged_text.strip():
                        translation_tasks.append((paragraph, tagged_text))
                 elif "text" in paragraph and paragraph["text"].strip():
                    translation_tasks.append((paragraph, paragraph["text"]))

        props_element = ET.SubElement(text_element, "properties")
        # Defer writing JSON to XML until after translation updates shape_data
        deferred_writes.append((props_element, shape_data))


def extract_text_from_slide(
    slide,
    slide_number: int,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
):
    """Extract text from a slide and optionally translate it using batching."""
    slide_element = ET.Element("slide")
    slide_element.set("number", str(slide_number))
    
    translation_tasks = [] # List of (paragraph_data, text_to_translate)
    deferred_writes = []   # List of (xml_element, data_object)
    
    for shape_index, shape in enumerate(slide.shapes):
        process_shape(
            shape, 
            str(shape_index), 
            slide_element, 
            translation_tasks,
            deferred_writes
        )
    
    # Batch Translate if we have tasks
    if translator and translation_tasks:
        texts_to_translate = [t[1] for t in translation_tasks]
        
        # Use batch JSON translation
        translated_texts = translator.translate_batch_json(texts_to_translate, source_lang, target_lang)
        
        # Apply results (Update the dictionary objects in memory)
        for (paragraph, original_text), translated_text in zip(translation_tasks, translated_texts):
            if "runs" in paragraph:
                parse_tagged_text_to_runs(translated_text, paragraph["runs"])
            else:
                paragraph["text"] = translated_text
    
    # Finally, write the UPDATED objects to XML
    for props_element, data_object in deferred_writes:
        props_element.text = json.dumps(data_object, indent=2, ensure_ascii=False)
            
    return slide_element


def ppt_to_xml(
    ppt_path: str,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
) -> Optional[str]:
    """Convert a PowerPoint presentation to XML."""
    root = ET.Element("presentation")
    base_dir = Path(ppt_path).parent
    try:
        prs = Presentation(ppt_path)
        root.set("file_path", Path(ppt_path).name)
        workers = max(1, max_workers)
        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_to_slide = {
                executor.submit(
                    extract_text_from_slide,
                    slide,
                    slide_number,
                    translator=translator,
                    source_lang=source_lang,
                    target_lang=target_lang,
                ): slide_number
                for slide_number, slide in enumerate(prs.slides, start=1)
            }
            # Sort results by slide number to ensure correct order in XML
            # (Though concurrent execution doesn't guarantee order, collecting futures does)
            
            # Wait for all futures
            results = []
            for future in future_to_slide:
                results.append((future_to_slide[future], future.result()))
            
            # Sort by slide number
            results.sort(key=lambda x: x[0])
            
            for slide_number, slide_element in results:
                root.append(slide_element)
                # We can skip saving individual slide XMLs to reduce IO unless debugging
                if False: # Disable per-slide debug file for performance/cleanup
                     intermediate_path = base_dir / f"slide_{slide_number}_{'translated' if translator else 'original'}.xml"
                     # ... saving logic ...
        
        return minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error processing presentation: {exc}")
        import traceback
        traceback.print_exc()
        return None


def process_shape_apply(shape, shape_index_str: str, xml_slide: ET.Element):
    """Recursively apply properties to a shape or group from XML."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if hasattr(shape, "shapes"):
            for child_idx, child_shape in enumerate(shape.shapes):
                child_index_str = f"{shape_index_str}:{child_idx}"
                process_shape_apply(child_shape, child_index_str, xml_slide)
        return

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table_element = xml_slide.find(f".//table_element[@shape_index='{shape_index_str}']")
        if table_element is not None:
            props_element = table_element.find("properties")
            if props_element is not None and props_element.text:
                try:
                    table_data = json.loads(props_element.text)
                    apply_table_properties(shape.table, table_data)
                except Exception as exc:  # pragma: no cover
                    print(f"Error applying table properties: {exc}")
    
    elif hasattr(shape, "text_frame"):
        text_element = xml_slide.find(f".//text_element[@shape_index='{shape_index_str}']")
        if text_element is not None:
            props_element = text_element.find("properties")
            if props_element is not None and props_element.text:
                try:
                    shape_data = json.loads(props_element.text)
                    apply_shape_properties(shape, shape_data)
                except Exception as exc:  # pragma: no cover
                    print(f"Error applying shape properties: {exc}")


def create_translated_ppt(original_ppt_path: str, translated_xml_path: str, output_ppt_path: str) -> None:
    """Create a new PowerPoint presentation using translated content."""
    try:
        prs = Presentation(original_ppt_path)
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        
        # Build map of slide number to xml element
        slide_map = {}
        for xml_slide in root.findall("slide"):
            slide_map[int(xml_slide.get("number"))] = xml_slide
            
        for slide_number, slide in enumerate(prs.slides, start=1):
            xml_slide = slide_map.get(slide_number)
            if xml_slide is None:
                continue
                
            for shape_index, shape in enumerate(slide.shapes):
                process_shape_apply(shape, str(shape_index), xml_slide)
                                
        prs.save(output_ppt_path)
        print(f"Translated PowerPoint saved to: {output_ppt_path}")
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Error creating translated PowerPoint: {exc}")
        import traceback
        traceback.print_exc()


def cleanup_intermediate_files(base_dir: Path, pattern: str = "slide_*.xml") -> None:
    """Remove intermediate XML files."""
    try:
        for file in base_dir.glob(pattern):
            file.unlink()
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Warning: Could not clean up intermediate files: {exc}")


def process_ppt_file(
    ppt_path: Path,
    *,
    translator: TranslationService,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
    cleanup: bool = True,
) -> Optional[Path]:
    """Process a single PowerPoint file from extraction to translated output."""
    if not ppt_path.is_file():
        raise FileNotFoundError(f"'{ppt_path}' is not a valid file.")
    if ppt_path.suffix.lower() not in {".ppt", ".pptx"}:
        raise ValueError(f"'{ppt_path}' is not a PowerPoint file.")

    base_dir = ppt_path.parent

    print(f"Generating original XML for {ppt_path.name}...")
    original_xml = ppt_to_xml(
        str(ppt_path),
        translator=None,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
    )
    if original_xml:
        original_output_path = base_dir / f"{ppt_path.stem}_original.xml"
        with open(original_output_path, "w", encoding="utf-8") as handle:
            handle.write(original_xml)
        print(f"Original XML saved: {original_output_path}")

    print(
        f"Generating translated XML (from {source_lang} to {target_lang}) for {ppt_path.name}..."
    )
    translated_xml = ppt_to_xml(
        str(ppt_path),
        translator=translator,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
    )
    if not translated_xml:
        return None

    translated_output_path = base_dir / f"{ppt_path.stem}_translated.xml"
    with open(translated_output_path, "w", encoding="utf-8") as handle:
        handle.write(translated_xml)
    print(f"Translated XML saved: {translated_output_path}")

    print(f"Creating translated PPT for {ppt_path.name}...")
    output_filename = f"{ppt_path.stem}_translated{ppt_path.suffix}"
    output_ppt_path = base_dir / output_filename
    create_translated_ppt(str(ppt_path), str(translated_output_path), str(output_ppt_path))

    if cleanup:
        cleanup_intermediate_files(base_dir)
        print("Cleanup complete.")

    return output_ppt_path
